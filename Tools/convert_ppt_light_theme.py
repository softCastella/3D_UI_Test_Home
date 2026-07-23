from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import glob, os

src = max(
    [p for p in glob.glob(os.path.join("Docs", "PPT", "*.pptx"))
     if not os.path.basename(p).startswith("~$") and "상세검토본" in p and "밝은테마" not in p],
    key=os.path.getmtime,
)
dst = os.path.join("Docs", "PPT", "화학물질_안전훈련_VR_사전기획_개발일정_상세검토본_밝은테마_v2.pptx")

prs = Presentation(src)

# Original dark-theme colors
NAVY = RGBColor(9, 20, 36)
NAVY2 = RGBColor(15, 32, 52)
CARD = RGBColor(22, 43, 66)
GRID = RGBColor(42, 66, 89)
WHITE = RGBColor(244, 248, 252)
MUTED = RGBColor(157, 177, 197)

# Light review/print palette
BG = RGBColor(247, 249, 252)
SURFACE = RGBColor(238, 244, 249)
PAPER = RGBColor(255, 255, 255)
BORDER = RGBColor(205, 217, 228)
INK = RGBColor(20, 38, 57)
SLATE = RGBColor(82, 103, 124)


def rgb_of(color):
    try:
        return color.rgb
    except (AttributeError, ValueError):
        return None


for slide in prs.slides:
    # Slide canvas
    fill = slide.background.fill
    if fill.type is not None:
        fill.solid()
        fill.fore_color.rgb = BG

    for shape in slide.shapes:
        # Shape fills: preserve cyan/blue/yellow/red accents, lighten dark surfaces.
        try:
            if shape.fill.type is not None:
                c = rgb_of(shape.fill.fore_color)
                if c == NAVY:
                    shape.fill.solid(); shape.fill.fore_color.rgb = BG
                elif c == NAVY2:
                    shape.fill.solid(); shape.fill.fore_color.rgb = SURFACE
                elif c == CARD:
                    shape.fill.solid(); shape.fill.fore_color.rgb = PAPER
        except (AttributeError, ValueError, TypeError):
            pass

        # Borders and separators
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                shape.line.fill.background()
            else:
                c = rgb_of(shape.line.color)
                if c in (NAVY, NAVY2, CARD, GRID):
                    shape.line.color.rgb = BORDER
        except (AttributeError, ValueError, TypeError):
            pass

        # Text colors
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                c = rgb_of(run.font.color)
                if c == WHITE:
                    run.font.color.rgb = INK
                elif c == MUTED:
                    run.font.color.rgb = SLATE

# The cover benefits from a slightly stronger white panel contrast.
cover = prs.slides[0]
cover.background.fill.solid(); cover.background.fill.fore_color.rgb = BG

prs.save(dst)
print(dst)
print("slides", len(prs.slides))
