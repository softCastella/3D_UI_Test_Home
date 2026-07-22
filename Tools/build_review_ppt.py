from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os

OUT = os.path.join("Docs", "PPT", "화학물질_안전훈련_VR_사전기획_개발일정_검토본.pptx")

W, H = Inches(13.333), Inches(7.5)
NAVY = RGBColor(9, 20, 36)
NAVY2 = RGBColor(15, 32, 52)
CYAN = RGBColor(37, 211, 198)
BLUE = RGBColor(55, 132, 255)
YELLOW = RGBColor(255, 193, 69)
RED = RGBColor(255, 97, 97)
WHITE = RGBColor(244, 248, 252)
MUTED = RGBColor(157, 177, 197)
CARD = RGBColor(22, 43, 66)
GRID = RGBColor(42, 66, 89)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
prs.core_properties.title = "화학물질 안전훈련 VR – 밀폐공간 안전교육"
prs.core_properties.subject = "사전 기획 및 개발 일정 검토본"
prs.core_properties.author = "3D_UI_Test_home 프로젝트"


def rect(slide, x, y, w, h, color, radius=False, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.color.rgb = line or color
    return shp


def text(slide, s, x, y, w, h, size=18, color=WHITE, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font="Malgun Gothic"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = s; r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box


def bullet_box(slide, items, x, y, w, h, size=16, color=WHITE, accent=CYAN, spacing=8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.level = 0; p.font.name = "Malgun Gothic"; p.font.size = Pt(size); p.font.color.rgb = color
        p.space_after = Pt(spacing); p._p.get_or_add_pPr().insert(0, p._p._new_buChar()) if False else None
        p.text = "•  " + item
    return box


def base(title, kicker="PROJECT REVIEW", section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = NAVY
    rect(slide, 0, 0, 13.333, .09, CYAN)
    text(slide, kicker, .55, .28, 4.5, .3, 9, CYAN, True)
    text(slide, title, .55, .68, 11.9, .65, 27, WHITE, True)
    rect(slide, .55, 1.43, 12.2, .02, GRID)
    if section:
        text(slide, section, 11.15, .30, 1.6, .3, 9, MUTED, True, PP_ALIGN.RIGHT)
    return slide


def footer(slide, n, source="사전 기획서 · 개발 일정"):
    text(slide, source, .55, 7.13, 5.5, .2, 8, MUTED)
    text(slide, f"{n:02d}", 12.15, 7.08, .6, .25, 9, MUTED, True, PP_ALIGN.RIGHT)


def card(slide, title_s, body, x, y, w, h, accent=CYAN, num=None, body_size=14):
    rect(slide, x, y, w, h, CARD, True, GRID)
    rect(slide, x, y, .06, h, accent, True, accent)
    if num:
        text(slide, num, x+.22, y+.18, .5, .3, 11, accent, True)
        tx = x+.72
    else: tx = x+.24
    text(slide, title_s, tx, y+.16, w-(tx-x)-.2, .35, 15, WHITE, True)
    if isinstance(body, list): bullet_box(slide, body, x+.24, y+.65, w-.48, h-.78, body_size, MUTED, spacing=5)
    else: text(slide, body, x+.24, y+.66, w-.48, h-.8, body_size, MUTED)


# 01 cover
s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
rect(s, 0, 0, 13.333, 7.5, NAVY)
for x, y, r in [(9.6,.7,2.7),(10.8,3.0,1.6),(8.2,4.7,2.0)]:
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(r), Inches(r)); shp.fill.background(); shp.line.color.rgb=GRID; shp.line.width=Pt(1.2)
rect(s, .68, .72, 1.25, .34, CYAN, True)
text(s, "VR SAFETY", .82, .78, 1.0, .18, 9, NAVY, True, PP_ALIGN.CENTER)
text(s, "화학물질 안전훈련 VR", .68, 1.55, 8.3, .65, 31, WHITE, True)
text(s, "밀폐공간 안전교육", .68, 2.27, 8.2, .78, 38, CYAN, True)
text(s, "사전 기획 · 개발 일정 통합 검토본", .72, 3.23, 6.3, .42, 18, MUTED)
rect(s, .72, 4.23, 5.9, 1.18, CARD, True, GRID)
text(s, "교육 목표", .98, 4.47, 1.1, .28, 12, CYAN, True)
text(s, "PPE 점검부터 작업허가 승인까지\n직접 수행하는 체험형 안전교육", 2.05, 4.40, 4.25, .7, 18, WHITE, True)
text(s, "Unity 6 · OpenXR · Meta Quest 2", .72, 6.65, 5.5, .3, 11, MUTED)
footer(s, 1)

# 02 agenda
s=base("검토 범위와 발표 구성", section="OVERVIEW")
items=[("01","기획 배경","교육 목적·대상·플레이어 역할"),("02","교육 설계","공통 조작과 2개 시나리오"),("03","체험 흐름","PPE → 현장 준비 → 승인"),("04","개발 계획","프로토타입·보완·QA·발표"),("05","관리 기준","완료 기준·범위 축소 원칙")]
for i,(n,t,b) in enumerate(items):
    y=1.75+i*.98; text(s,n,.72,y,.55,.35,13,CYAN,True); text(s,t,1.45,y,2.0,.35,17,WHITE,True); text(s,b,3.5,y+.01,7.8,.35,15,MUTED); rect(s,11.95,y+.05,.25,.25,CYAN,True)
footer(s,2)

# 03 problem
s=base("왜 VR 밀폐공간 안전교육인가", section="01 PLANNING")
card(s,"기존 교육의 한계",["설명·문서 중심 교육은 절차 기억과 행동 전환에 한계","실제 설비에서 반복 훈련하기 어렵고 사고 위험이 큼","PPE·LOTO·가스 측정·환기의 순서 오류를 체감하기 어려움"],.65,1.75,3.85,4.7,RED,body_size=15)
card(s,"VR이 제공하는 가치",["위험 없이 동일 절차를 반복 수행","오답 즉시 피드백과 재시도로 행동 교정","설비·장비·공간 맥락 안에서 순서 학습","Quest 기반 독립 실행으로 교육 접근성 확보"],4.75,1.75,3.85,4.7,CYAN,body_size=15)
card(s,"프로젝트의 해법","혼합기 내부 청소 전 과정을 하나의 작업 시나리오로 연결하고, 플레이어가 직접 확인·점검·차단·측정·설치·보고하도록 설계한다.",8.85,1.75,3.85,4.7,BLUE,body_size=17)
footer(s,3)

# 04 overview
s=base("프로젝트 개요", section="01 PLANNING")
card(s,"프로젝트명","화학물질 안전훈련 VR\n– 밀폐공간 안전교육",.65,1.75,3.8,1.55,CYAN,body_size=17)
card(s,"교육 대상",["화학제품 제조공장 생산작업자","혼합기 내부 청소 지정 작업자","신규·정기 안전교육 대상자"],4.75,1.75,3.8,2.25,BLUE)
card(s,"현장 설정",["중소규모 화학제품 제조공장","2층 작업대 높이의 혼합기","계단·작업대를 통한 맨홀 접근"],8.85,1.75,3.8,2.25,YELLOW)
card(s,"핵심 수행 범위",["PPE 확인 및 착용","구동 전원 차단·LOTO","1·2차 가스 측정과 환기","통신 확인·작업허가 승인"],.65,4.3,12.0,1.75,CYAN,body_size=16)
footer(s,4)

# 05 roles
s=base("플레이어 역할과 교육 목표", section="01 PLANNING")
card(s,"플레이어","혼합기 내부 청소작업을 수행하도록 지정된 작업자",.65,1.85,3.65,1.55,CYAN,body_size=16)
card(s,"감시인","현장 외부에서 준비상태·통신상태를 확인하고 승인 결과 전달",4.85,1.85,3.65,1.55,BLUE,body_size=16)
card(s,"교육 시스템","현재 단계 안내, 조작 유도, 완료·오류 피드백, 진행 상태 관리",9.05,1.85,3.65,1.55,YELLOW,body_size=16)
text(s,"학습 성과",.65,4.0,2,.35,18,CYAN,True)
for i,t in enumerate(["필요 PPE를 식별하고 정상·불량을 판정한다","밀폐공간 진입 전 준비절차를 올바른 순서로 수행한다","모든 조건을 완료한 뒤 작업허가를 받는 원칙을 이해한다"]):
    rect(s,.65,4.55+i*.62,.34,.34,CYAN,True); text(s,str(i+1),.65,4.60+i*.62,.34,.18,10,NAVY,True,PP_ALIGN.CENTER); text(s,t,1.15,4.52+i*.62,10.8,.4,16,WHITE)
footer(s,5)

# 06 structure
s=base("교육 콘텐츠 구조", section="02 DESIGN")
text(s,"PPE룸",.75,1.75,2,.35,15,CYAN,True)
stages=[("시작 안내","조작 이해"),("시나리오 선택","교육 선택"),("PPE 착용 교육","점검·판정·착용"),("현장 이동","씬 전환"),("준비작업 교육","차단·측정·환기"),("작업허가 승인","완료·복귀")]
for i,(a,b) in enumerate(stages):
    x=.65+i*2.08; col=CYAN if i<3 else BLUE
    rect(s,x,2.35,1.78,1.42,CARD,True,col); text(s,f"{i+1:02d}",x+.16,2.53,.4,.25,11,col,True); text(s,a,x+.16,2.89,1.45,.4,14,WHITE,True,PP_ALIGN.CENTER); text(s,b,x+.16,3.34,1.45,.26,10,MUTED,False,PP_ALIGN.CENTER)
    if i<5: text(s,"→",x+1.8,2.87,.28,.35,18,MUTED,True,PP_ALIGN.CENTER)
text(s,"두 시나리오는 독립 선택이 가능하지만, 최종 체험은 ‘개인 준비 → 현장 준비 → 승인’의 하나의 안전 행동 흐름으로 연결된다.",.9,4.55,11.5,.85,19,WHITE,True,PP_ALIGN.CENTER)
rect(s,2.4,5.7,8.55,.65,NAVY2,True,GRID); text(s,"설명형 콘텐츠가 아닌 직접 수행형 안전교육",2.65,5.9,8.05,.25,16,CYAN,True,PP_ALIGN.CENTER)
footer(s,6)

# 07 common UX
s=base("공통 진행·조작 규칙", section="02 DESIGN")
controls=[("TRIGGER","레이 선택 · 텔레포트 · 측정 · 송신"),("GRIP","한 번 눌러 잡기 / 다시 눌러 놓기"),("B","확인 · 다음 페이지 · 대화 종료 · 장비 작동"),("A","PPE 교육에서 불량 PPE 폐기")]
for i,(k,v) in enumerate(controls):
    x=.65+(i%2)*6.15; y=1.75+(i//2)*1.25
    rect(s,x,y,5.85,1.0,CARD,True,GRID); rect(s,x+.18,y+.18,1.15,.64,CYAN if i<2 else BLUE,True); text(s,k,x+.18,y+.37,1.15,.22,12,NAVY,True,PP_ALIGN.CENTER); text(s,v,x+1.55,y+.31,4.0,.35,15,WHITE,True)
card(s,"안내 시스템",["음성 안내방송 + 화면 안내창","큰 원형 위치 마커 / 작은 원형 아이템 마커","완료는 초록 체크, 오류는 빨간 X와 재안내"],.65,4.55,12.0,1.55,YELLOW,body_size=15)
footer(s,7)

# 08 PPE flow
s=base("시나리오 A | PPE 착용 안전교육", section="03 EXPERIENCE")
ppe=["작업계획서\n확인","PPE 단스\n이동","표면 상태\n육안점검","사용·폐기\n판정","착용 상태\n저장","거울 확인","현장 이동"]
for i,t in enumerate(ppe):
    x=.55+i*1.78; rect(s,x,2.2,1.48,1.25,CARD,True,CYAN); text(s,str(i+1),x+.15,2.38,.3,.25,11,CYAN,True); text(s,t,x+.2,2.72,1.08,.52,14,WHITE,True,PP_ALIGN.CENTER)
    if i<6:text(s,"›",x+1.49,2.59,.28,.4,22,MUTED,True,PP_ALIGN.CENTER)
text(s,"필요 PPE",.7,4.25,1.5,.3,15,CYAN,True)
for i,name in enumerate(["내화학 방호복","내화학 장화","내화학 장갑","송기마스크","안전모"]):
    x=.7+i*2.42; rect(s,x,4.85,2.1,.75,NAVY2,True,GRID); text(s,name,x+.1,5.09,1.9,.26,14,WHITE,True,PP_ALIGN.CENTER)
text(s,"정상 PPE만 착용하고, 불량 PPE는 폐기하는 판정 행동을 핵심 학습으로 구성",.75,6.18,11.8,.4,16,YELLOW,True,PP_ALIGN.CENTER)
footer(s,8)

# 09 PPE interactions
s=base("PPE 교육 | 핵심 인터랙션과 피드백", section="03 EXPERIENCE")
card(s,"작업계획서",["태블릿을 그립으로 잡기","작업내용·필요 PPE 체크","전체 확인 후 다음 단계 활성화"],.65,1.75,3.75,2.05,CYAN)
card(s,"육안점검",["PPE를 손에 들고 회전 확인","찢어짐·균열·오염·변색 표현","정상/불량에 따라 사용 또는 폐기"],4.8,1.75,3.75,2.05,BLUE)
card(s,"착용·거울",["PPE가 몸 방향으로 이동 후 착용 상태 저장","신체 대신 착용 PPE만 표시","거울에서 연결부·표면·안전모 확인"],8.95,1.75,3.75,2.05,YELLOW)
rect(s,.65,4.25,5.75,1.4,CARD,True,GRID); text(s,"✓  올바른 선택",.95,4.56,2.2,.32,18,CYAN,True); text(s,"초록 체크 · 완료 SFX · 확인멘트",.95,5.03,4.7,.3,14,MUTED)
rect(s,6.9,4.25,5.8,1.4,CARD,True,GRID); text(s,"×  잘못된 선택",7.2,4.56,2.2,.32,18,RED,True); text(s,"빨간 X · 오류 SFX · 재확인 안내",7.2,5.03,4.7,.3,14,MUTED)
footer(s,9)

# 10 field flow
s=base("시나리오 B | 밀폐공간 진입 전 준비작업", section="03 EXPERIENCE")
steps=["감시인\n최초 대화","안전 문서\n확인","전원 OFF","LOTO","1차 가스\n측정","환기팬·덕트\n설치","환기","2차 가스\n측정","무전 테스트","승인·복귀"]
for i,t in enumerate(steps):
    row=i//5; col=i%5; x=.7+col*2.48; y=1.8+row*1.58
    rect(s,x,y,2.06,1.13,CARD,True,BLUE if row else CYAN); text(s,f"{i+1:02d}",x+.15,y+.15,.38,.22,10,CYAN if row==0 else BLUE,True); text(s,t,x+.25,y+.48,1.56,.46,14,WHITE,True,PP_ALIGN.CENTER)
    if col<4:text(s,"→",x+2.08,y+.41,.35,.35,16,MUTED,True,PP_ALIGN.CENTER)
text(s,"승인 조건",.7,5.35,1.2,.3,14,YELLOW,True)
text(s,"문서 확인 + 전원 차단 + LOTO + 1·2차 가스 측정 + 환기 + 무전 테스트 완료",1.9,5.31,10.7,.45,17,WHITE,True)
footer(s,10)

# 11 safety docs/loto
s=base("현장 준비 ① | 문서 확인 · 전원 차단 · LOTO", section="03 EXPERIENCE")
card(s,"안전 문서 확인",["작업계획서","밀폐공간 작업허가서","MSDS · 작업 절차 · 비상대응 정보","마지막 페이지까지 확인해야 완료"],.65,1.75,3.75,3.85,CYAN,body_size=15)
card(s,"혼합기 전원 OFF",["혼합기 구동 전원 상태 확인","B 버튼으로 OFF 처리","표시등 소등과 상태값 일치","OFF 이후 재가동 방지"],4.8,1.75,3.75,3.85,BLUE,body_size=15)
card(s,"LOTO 설치",["전원 OFF 후 자물쇠·태그 활성화","그립 토글로 잡아 지정 위치 스냅","설치 후 재조작 차단","완료 시 초록 체크"],8.95,1.75,3.75,3.85,YELLOW,body_size=15)
text(s,"절차 순서 강제: 문서 확인 → 전원 OFF → LOTO",.7,6.05,12,.4,18,CYAN,True,PP_ALIGN.CENTER)
footer(s,11)

# 12 gas
s=base("현장 준비 ② | 1·2차 가스 측정", section="03 EXPERIENCE")
for i,(n,t,b) in enumerate([("01","측정기 잡기","오른손 그립 토글"),("02","측정 시작","맨홀 방향에서 트리거"),("03","자동 연출","왼손은 외부, 호스·프로브만 내부"),("04","구간 측정","상부·중부·하부 순차 표시"),("05","결과 확인","5초 게이지 후 정상수치 판정")]):
    x=.65+i*2.45; rect(s,x,1.95,2.05,2.05,CARD,True,CYAN if i<3 else BLUE); text(s,n,x+.18,2.16,.4,.24,10,CYAN,True); text(s,t,x+.18,2.63,1.7,.35,16,WHITE,True); text(s,b,x+.18,3.12,1.7,.6,13,MUTED)
text(s,"교육 핵심",.75,4.65,1.2,.3,14,YELLOW,True)
text(s,"환기 전 1차 측정과 환기 후 2차 측정을 구분하고, 한 번의 조작으로 상·중·하부 측정 의미를 전달한다.",.75,5.12,11.75,.75,19,WHITE,True,PP_ALIGN.CENTER)
footer(s,12)

# 13 ventilation radio approval
s=base("현장 준비 ③ | 환기 · 통신 · 작업허가", section="03 EXPERIENCE")
card(s,"환기팬·덕트",["팬 쪽 덕트 스냅 후 전원 작동","반대쪽 덕트를 맨홀에 설치","10초 교육용 환기 게이지","준비작업 종료까지 팬 가동"],.65,1.75,3.75,3.85,CYAN,body_size=15)
card(s,"무전 테스트",["무전기를 잡은 상태에서만 송신","플레이어 → 잡음 → 감시인 순서","중복 입력·음성 겹침 방지","완료 후 승인 단계 활성화"],4.8,1.75,3.75,3.85,BLUE,body_size=15)
card(s,"작업허가 승인",["모든 준비 조건 완료 여부 확인","감시인 최종 대화","작업허가서 승인 도장","성공 메시지 후 PPE룸 복귀"],8.95,1.75,3.75,3.85,YELLOW,body_size=15)
footer(s,13)

# 14 UX assets
s=base("UI · SFX · VFX 설계", section="03 EXPERIENCE")
card(s,"UI",["World Space 안내창·대화창","태블릿 문서·체크박스","가스 측정·환기 게이지","성공 메시지·복귀 버튼"],.65,1.75,3.75,4.25,CYAN,body_size=15)
card(s,"SFX / VO",["버튼·잡기·체결·측정 완료음","환기팬 3D 공간음","무전 송수신·잡음","감시인 안내와 완료 음성"],4.8,1.75,3.75,4.25,BLUE,body_size=15)
card(s,"VFX / MOTION",["위치·아이템 마커 반복 모션","대상 하이라이트·배경 블러","PPE 손상·착용 페이드","팬 회전·승인 도장·씬 페이드"],8.95,1.75,3.75,4.25,YELLOW,body_size=15)
text(s,"원칙: 시야를 가리지 않고, 현재 단계의 정보만 제공하며, 음성·화면 지시가 일치해야 한다.",.75,6.35,11.75,.4,16,WHITE,True,PP_ALIGN.CENTER)
footer(s,14)

# 15 schedule summary
s=base("전체 개발 로드맵", section="04 SCHEDULE")
phases=[("07.20–08.09","프로토타입 제작","핵심 시나리오 전체 흐름"),("08.10","프로토타입 시연","작동 가능한 수직형 빌드"),("08.11–08.23","기능 완성·보완","피드백·연출·에셋 적용"),("08.24–09.06","통합·최적화","Quest 2·문서·안정 빌드"),("09.07–09.13","QA 주간","기능·조작·절차 집중 검수"),("09.14–09.18","리허설·제출","발표·시연·최종 평가")]
for i,(d,t,b) in enumerate(phases):
    y=1.65+i*.84; col=CYAN if i<2 else BLUE if i<4 else YELLOW
    rect(s,.75,y,.16,.62,col,True); text(s,d,1.15,y+.08,1.75,.26,12,col,True); text(s,t,3.0,y+.03,2.25,.3,16,WHITE,True); text(s,b,5.55,y+.05,6.4,.3,14,MUTED)
footer(s,15)

# 16 weekly prototype
s=base("프로토타입 3주 실행 계획", section="04 SCHEDULE")
card(s,"1주차 · 07.20–07.26",["공통 진행 시스템","레이 UI·위치/아이템 마커","그립 토글·A/B/트리거","안내창·정답/오답 피드백"],.65,1.75,3.75,3.95,CYAN,body_size=15)
card(s,"2주차 · 07.27–08.02",["작업계획서 태블릿","PPE 점검·사용·폐기","착용 상태 저장","거울 확인·현장 이동 연결"],4.8,1.75,3.75,3.95,BLUE,body_size=15)
card(s,"3주차 · 08.03–08.09",["감시인·문서·전원·LOTO","1·2차 가스 측정","환기팬·덕트·무전","승인·복귀·Quest 2 안정화"],8.95,1.75,3.75,3.95,YELLOW,body_size=15)
text(s,"08.09 신규 기능 추가 금지 → 시연 순서·초기화·입력·빌드 안정성만 점검",.7,6.15,12,.4,17,RED,True,PP_ALIGN.CENTER)
footer(s,16)

# 17 demo scope
s=base("8월 10일 프로토타입 시연 기준", section="04 SCHEDULE")
card(s,"시연 범위",["시나리오 선택 → PPE 교육","현장 문서 → OFF → LOTO","가스 측정 → 환기 → 재측정","무전 → 감시인 승인 → 성공"],.65,1.75,3.75,4.25,CYAN,body_size=15)
card(s,"허용되는 임시 요소",["임시 안내 음성·문서 이미지","단순화된 자동 손 애니메이션","일부 미완성 VFX·SFX","기본 형태 성공 메시지"],4.8,1.75,3.75,4.25,BLUE,body_size=15)
card(s,"반드시 만족",["처음부터 끝까지 중단 없음","다음 행동을 스스로 이해","조작과 안내 내용 일치","에디터 조작 없이 시연 가능"],8.95,1.75,3.75,4.25,YELLOW,body_size=15)
footer(s,17)

# 18 post prototype
s=base("프로토타입 이후 완성 전략", section="04 SCHEDULE")
ph=[("08.11–08.16","1차 보완","진행 오류·조작 설명·UI·동선·대화 수정"),("08.17–08.23","2차 완성","PPE 셰이더·거울·자동 손·호스·SFX·승인 연출"),("08.24–08.30","통합","전체 연결·초기화·복구·중복 방지·전체 플레이"),("08.31–09.06","최적화·문서","Quest 2 성능·오디오·텍스처·문서·스크린샷")]
for i,(d,t,b) in enumerate(ph):
    y=1.75+i*1.1; rect(s,.7,y,2.2,.82,NAVY2,True,GRID); text(s,d,.86,y+.13,1.85,.25,13,CYAN,True); text(s,t,.86,y+.46,1.85,.22,12,WHITE,True)
    rect(s,3.15,y,9.45,.82,CARD,True,GRID); text(s,b,3.45,y+.25,8.85,.35,15,WHITE,True)
text(s,"우선순위: 시나리오 안정성 → 조작 명확성 → Quest 성능 → 시청각 완성도",.75,6.45,11.8,.35,16,YELLOW,True,PP_ALIGN.CENTER)
footer(s,18)

# 19 QA
s=base("QA · 리허설 · 제출 계획", section="04 SCHEDULE")
card(s,"QA · 09.07–09.13",["신규 기능·에셋 추가 금지","진행 불가 → 입력 → 초기화 → 절차 오류 순 수정","기능별 QA 후 전체 반복 플레이","Quest 2 실기 빌드·회귀 테스트"],.65,1.75,3.75,4.35,CYAN,body_size=15)
card(s,"리허설 · 09.14–09.16",["발표 대본·예상 질문 정리","VR 시연 동선과 시간 측정","화면 공유·대체 영상 점검","발표 전날 비치명적 수정 금지"],4.8,1.75,3.75,4.35,BLUE,body_size=15)
card(s,"제출 · 09.17–09.18",["VR 실행 빌드·기획/설계/R&D 문서","개발 일정·QA 결과·발표자료","시연 또는 백업 영상·에셋 출처","최종 백업과 과정 평가"],8.95,1.75,3.75,4.35,YELLOW,body_size=15)
footer(s,19)

# 20 scope control
s=base("일정 지연 시 범위 조정 원칙", section="05 CONTROL")
card(s,"반드시 유지",["시나리오 진행·입력·마커","문서 확인·PPE 판정","전원 OFF·LOTO·가스 측정","환기·무전·승인·복귀"],.65,1.75,3.75,4.45,CYAN,body_size=15)
card(s,"단순화 가능",["PPE 손상 셰이더 정밀도","거울 해상도·자동 손 세부 동작","호스 변형·덕트 물리","공기 흐름 VFX·NPC 애니메이션"],4.8,1.75,3.75,4.45,YELLOW,body_size=15)
card(s,"삭제 가능",["장식용 환경 애니메이션","진행과 무관한 오브젝트 조작","불필요한 자유 이동","복잡한 컷신·추가 UI"],8.95,1.75,3.75,4.45,RED,body_size=15)
footer(s,20)

# 21 criteria
s=base("완료 기준과 검토 포인트", section="05 CONTROL")
left=["두 교육 시나리오가 처음부터 끝까지 실행되는가","플레이어가 현재·다음 행동을 이해하는가","절차 순서를 건너뛸 수 없도록 제어되는가","입력 중복·오브젝트 유실에서 복구 가능한가"]
right=["PPE 판정과 착용 결과가 명확한가","가스 측정·환기·통신의 교육 의미가 전달되는가","모든 조건 완료 후에만 승인이 실행되는가","Meta Quest 2에서 안정적으로 작동하는가"]
card(s,"콘텐츠·진행",left,.65,1.75,5.85,4.55,CYAN,body_size=16)
card(s,"교육·기술",right,6.8,1.75,5.85,4.55,BLUE,body_size=16)
rect(s,2.45,6.52,8.4,.5,CYAN,True); text(s,"검토 의견을 반영해 R&D 상세본과 최종 발표본으로 확장",2.65,6.65,8.0,.2,13,NAVY,True,PP_ALIGN.CENTER)
footer(s,21)

# 22 PPE interaction specification
s=base("PPE 시나리오 | 단계별 인터랙션 명세", section="03 INTERACTION")
headers=[("단계",.65,1.45),("플레이어 행동 / 입력",2.65,3.35),("완료 조건",6.15,2.45),("오류·복구 / 출력",8.75,3.9)]
for h,x,w in headers: rect(s,x,1.62,w,.48,CYAN,True); text(s,h,x+.08,1.76,w-.16,.2,11,NAVY,True,PP_ALIGN.CENTER)
rows=[
 ("작업계획서","태블릿 그립 → 체크박스 트리거","전 항목 확인","미확인 시 진행 차단 / 초록 체크"),
 ("PPE 이동","위치 마커 레이 + 트리거","지정 위치·방향 도착","오배치 방지 / 대상 PPE 마커 활성"),
 ("육안점검","PPE 그립 토글 → 손목 회전","표면 상태 확인","놓치면 원위치 복구 / 선택 패널"),
 ("사용·폐기","정상: 트리거, 불량: A → B 확정","상태와 선택 일치","오답 빨간 X / 재시도·상태 초기화"),
 ("착용 처리","정상 PPE 사용 확정","착용 상태 저장","잡기 해제 / 몸 방향 이동·페이드"),
 ("거울 확인","위치 이동 → 안내 항목 시선 확인","5종 PPE 확인 완료","거울 종료 / 현장 이동 마커 활성")]
for i,row in enumerate(rows):
 y=2.18+i*.72; fill=NAVY2 if i%2==0 else CARD
 for val,x,w in zip(row,[.65,2.65,6.15,8.75],[2.0,3.5,2.6,3.9]):
  rect(s,x,y,w,.62,fill,False,GRID); text(s,val,x+.11,y+.13,w-.22,.34,11,WHITE if x<8 else MUTED,x==.65,PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
footer(s,22,"사전 기획서 4–5장")

# 23 field interaction specification A
s=base("현장 시나리오 | 문서·전원·LOTO·가스 측정", section="03 INTERACTION")
items=[
 ("안전 문서","태블릿 그립 / B로 페이지 이동","마지막 페이지 확인","중복 입력 잠금, 배경 블러 종료"),
 ("전원 OFF","전원 상태 확인 / B","표시등·내부 상태 OFF","LOTO 선행 차단, 다시 ON 불가"),
 ("LOTO","자물쇠·태그 그립 → 스냅","지정 위치 설치 완료","오스냅 차단, 설치 후 잡기 비활성"),
 ("가스측정기","오른손 그립 / 맨홀에서 트리거","측정 조건 4개 충족","거리·단계 확인, 중복 측정 잠금"),
 ("자동 왼손·호스","측정 시작 시 자동 실행","호스·프로브만 내부 삽입","텔레포트 잠금, 종료 후 손 복원"),
 ("상·중·하부 UI","5초 게이지 순차 표시","종합 결과 확인","1·2차 결과 구분, 정상 시 체크")]
for i,(a,b,c,d) in enumerate(items):
 y=1.62+i*.82; col=CYAN if i<3 else BLUE
 rect(s,.65,y,1.72,.7,CARD,True,col); text(s,a,.82,y+.2,1.4,.28,13,WHITE,True,PP_ALIGN.CENTER)
 text(s,b,2.62,y+.09,3.15,.5,12,WHITE); text(s,c,5.95,y+.09,2.55,.5,12,CYAN); text(s,d,8.72,y+.09,3.7,.5,12,MUTED)
text(s,"행동 / 입력",2.62,6.63,3,.2,10,MUTED); text(s,"완료 조건",5.95,6.63,2,.2,10,CYAN); text(s,"오류 방지·복구",8.72,6.63,3,.2,10,MUTED)
footer(s,23,"사전 기획서 6장")

# 24 field interaction specification B
s=base("현장 시나리오 | 환기·무전·승인", section="03 INTERACTION")
items=[
 ("덕트 팬 연결","덕트 끝 그립 → 팬 연결부 스냅","팬 쪽 연결 완료","맨홀 쪽 선행 설치 차단"),
 ("환기팬 작동","연결 후 B 버튼","팬 회전·SFX ON","반복 B 무시, 승인까지 ON 유지"),
 ("덕트 맨홀 연결","반대 끝 그립 → 이동 → 스냅","양단 연결 완료","단계별 모델 교체로 물리 오류 축소"),
 ("환기 게이지","양단 연결 후 자동 시작","10초 완료","진행 중 다음 단계 차단 / 2차 측정 활성"),
 ("무전 테스트","무전기 그립 → 트리거 송신","감시인 응답까지 재생","중복 송신 잠금, 음성 순차 재생"),
 ("작업허가 승인","감시인 레이 선택","7개 승인 조건 충족","미완료 항목 안내 / 도장 1회 실행")]
for i,(a,b,c,d) in enumerate(items):
 y=1.62+i*.82; col=YELLOW if i<4 else CYAN
 rect(s,.65,y,1.72,.7,CARD,True,col); text(s,a,.82,y+.2,1.4,.28,13,WHITE,True,PP_ALIGN.CENTER)
 text(s,b,2.62,y+.09,3.15,.5,12,WHITE); text(s,c,5.95,y+.09,2.55,.5,12,col); text(s,d,8.72,y+.09,3.7,.5,12,MUTED)
footer(s,24,"사전 기획서 6장")

# 25 week 1 daily
s=base("1주차 상세 일정 | 공통 시스템 구축", section="04 DAILY PLAN")
days=[("07.20 월","백업·씬 구조 확인·범위 고정·단계 목록·기존 코드 점검"),("07.21 화","카드 연결·상세 오버레이·교육 선택·돌아가기·레이 UI 검증"),("07.22 수","위치 마커 프리팹·알파 모션·고정 텔레포트·방향·페이드"),("07.23 목","아이템 마커 프리팹·크기 모션·단계별 활성·오브젝트 연동"),("07.24 금","그립 토글·트리거·A/B 문맥 기능·씬별 입력 분리"),("07.25–26","Quest 입력·마커 테스트·UI 오류·안내창·피드백·SFX·백업")]
for i,(d,b) in enumerate(days):
 y=1.55+i*.78; rect(s,.72,y,1.55,.61,CYAN if i<5 else YELLOW,True); text(s,d,.82,y+.18,1.35,.22,11,NAVY,True,PP_ALIGN.CENTER)
 rect(s,2.5,y,10.0,.61,CARD,True,GRID); text(s,b,2.75,y+.16,9.5,.3,13,WHITE,True)
text(s,"주차 완료 기준  |  시나리오 선택 · 고정 이동 · 현재 대상 마커 · 안정적 그립 토글 · 문맥별 버튼 · 완료/오류 피드백",.78,6.38,11.7,.48,14,CYAN,True,PP_ALIGN.CENTER)
footer(s,25,"개발 일정 9-2")

# 26 week 2 daily
s=base("2주차 상세 일정 | PPE 착용 교육", section="04 DAILY PLAN")
days=[("07.27 월","PPE 시작 단계·작업계획서 태블릿·아이템 마커·그립 토글"),("07.28 화","체크박스·체크 애니메이션·전체 확인 판정·오버레이 종료"),("07.29 수","PPE 단스 이동·PPE별 마커·잡기/놓기·손안 회전 확인"),("07.30 목","사용/폐기/취소 패널·정상/불량 데이터·A/B 판정·오답 복구"),("07.31 금","몸 방향 이동·착용 상태 저장·착용 모델·전체 완료 판정"),("08.01–02","거울 RT·레이어·성능 검증·전체 플레이·오류 수정·현장 연결")]
for i,(d,b) in enumerate(days):
 y=1.55+i*.78; rect(s,.72,y,1.55,.61,BLUE if i<5 else YELLOW,True); text(s,d,.82,y+.18,1.35,.22,11,NAVY,True,PP_ALIGN.CENTER)
 rect(s,2.5,y,10.0,.61,CARD,True,GRID); text(s,b,2.75,y+.16,9.5,.3,13,WHITE,True)
text(s,"주차 완료 기준  |  문서 확인 → PPE 육안점검 → 정상/불량 판정 → 착용 저장 → 거울 확인 → 현장 이동",.78,6.38,11.7,.48,14,BLUE,True,PP_ALIGN.CENTER)
footer(s,26,"개발 일정 9-3")

# 27 week 3 daily
s=base("3주차 상세 일정 | 밀폐공간 준비작업", section="04 DAILY PLAN")
days=[("08.03 월","현장 진입·감시인 NPC·최초 대화·목표 안내·문서 태블릿"),("08.04 화","전원 OFF·표시등·LOTO 자물쇠/태그·스냅·완료 판정"),("08.05 수","측정기·트리거·5초 게이지·상중하 UI·1차 측정·호스 우선"),("08.06 목","환기팬 이동·덕트 양단 연결·B 작동·SFX·10초 환기"),("08.07 금","2차 측정·무전 송수신·감시인 보고·승인 도장·복귀"),("08.08–09","신규 기능 중단·통합/입력/씬 전환·Quest 빌드·시연 2회·백업")]
for i,(d,b) in enumerate(days):
 y=1.55+i*.78; rect(s,.72,y,1.55,.61,YELLOW if i<5 else RED,True); text(s,d,.82,y+.18,1.35,.22,11,NAVY,True,PP_ALIGN.CENTER)
 rect(s,2.5,y,10.0,.61,CARD,True,GRID); text(s,b,2.75,y+.16,9.5,.3,13,WHITE,True)
text(s,"08.07 핵심 기능 완료 → 08.08 전체 통합 → 08.09 시연 빌드 확정 → 08.10 프로토타입 시연",.78,6.38,11.7,.48,14,RED,True,PP_ALIGN.CENTER)
footer(s,27,"개발 일정 9-4 · 10")

# 28 post-prototype weekly detail
s=base("시연 이후 주차별 산출물과 마감", section="04 DELIVERY PLAN")
rows=[("08.11–16","피드백 반영","진행·입력·UI·동선·대화 오류 수정","수정 빌드 + 피드백 처리표"),("08.17–23","기능 완성","셰이더·거울·자동 손·호스·덕트·SFX·도장","기능 완성 후보 빌드"),("08.24–30","통합","전체 연결·초기화·오브젝트 복구·SFX/VFX","통합 빌드 + 오류 목록"),("08.31–09.06","최적화·문서","거울/그림자/텍스처/오디오·기획/R&D/일정","QA 후보 빌드 + 문서 초안"),("09.07–13","QA","기능별 QA·전체 반복·회귀·Quest 실기","최종 후보 빌드 + QA 결과"),("09.14–18","리허설·제출","대본·동선·시간·화면공유·백업 영상·최종 제출","발표본 + 실행 빌드 + 제출 패키지")]
for i,(d,p,b,o) in enumerate(rows):
 y=1.54+i*.82; col=CYAN if i<2 else BLUE if i<4 else YELLOW
 rect(s,.65,y,1.5,.68,col,True); text(s,d,.74,y+.21,1.32,.22,10,NAVY,True,PP_ALIGN.CENTER)
 text(s,p,2.4,y+.17,1.42,.3,13,col,True); text(s,b,3.95,y+.1,5.25,.48,12,WHITE); text(s,o,9.45,y+.1,3.15,.48,12,MUTED)
text(s,"기간",.72,6.58,1,.2,9,MUTED); text(s,"목표",2.4,6.58,1,.2,9,MUTED); text(s,"핵심 작업",3.95,6.58,1,.2,9,MUTED); text(s,"검토 산출물",9.45,6.58,1,.2,9,MUTED)
footer(s,28,"개발 일정 11–15")

# Keep a review-focused 25-slide order: remove broad introductory repetition and
# place interaction/daily-plan specifications next to their summaries.
order=[1,2,4,6,7,8,22,9,10,23,24,11,12,13,15,16,25,26,27,17,18,28,19,20,21]
sld_ids=list(prs.slides._sldIdLst)
excluded=set(range(1,len(sld_ids)+1))-set(order)
for idx in excluded:
    prs.part.drop_rel(sld_ids[idx-1].rId)
for node in list(prs.slides._sldIdLst): prs.slides._sldIdLst.remove(node)
for idx in order: prs.slides._sldIdLst.append(sld_ids[idx-1])

OUT = os.path.join("Docs", "PPT", "화학물질_안전훈련_VR_사전기획_개발일정_상세검토본_정리.pptx")
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(OUT)
print("slides", len(prs.slides))
